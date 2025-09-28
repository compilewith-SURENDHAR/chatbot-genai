# backend/services/RAG/text_extractor.py

import os
import PyPDF2
import docx
from pptx import Presentation
import openpyxl


def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text


def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
    return text


def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def GetChunks(text, chunk_size=500, overlap=100):
    chunks = []

    step = chunk_size - overlap
    for i in range(0, len(text), step):
        chunk = text[i:i + chunk_size]
        chunks.append(chunk)

    return chunks

x=extract_text(r"C:\Users\surendhar\Downloads\Website Testing.pptx")
#print(x)  
#chunks = GetChunks(x)
#print("\n\n",chunks, "\n\n", len(chunks))
#print("\n \n", chunks[0])
